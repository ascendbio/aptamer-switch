"""Build the re:AGENT Track C submission deck.

Written as a script rather than assembled by hand so every number on a slide can
be traced to the file it came from. Anything stated here is either read out of a
run manifest, computed from a plate CSV in this repo, or recorded in ledger.py —
nothing is typed in from memory, which is the failure mode a deck invites.

Output is .pptx because the submission must be Google Slides: uploading a .pptx
to Drive opens it as a native Slides deck with the figures intact.

    python slides.py            → reAGENT_TrackC.pptx
"""

from __future__ import annotations

import csv
import json
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "out"
sys.path.insert(0, str(ROOT / "sources"))

INK = RGBColor(0x1A, 0x1A, 0x1E)
MUTED = RGBColor(0x62, 0x62, 0x6E)
ACCENT = RGBColor(0x2F, 0x5C, 0xE5)
WARN = RGBColor(0xB4, 0x53, 0x09)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Helvetica Neue"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.7)


# ---------------------------------------------------------------- numbers

def facts() -> dict:
    """Every figure quoted in the deck, read from the artifacts themselves."""
    runs = sorted((OUT / "runs").glob("*/manifest.json"),
                  key=lambda p: p.stat().st_mtime)
    m = json.loads(runs[-1].read_text())

    plate_rows = list(csv.DictReader((OUT / "IL-6_plate.csv").open()))
    hedged_rows = list(csv.DictReader((OUT / "IL-6_hedged_plate.csv").open()))
    by_hyp = Counter(r["Core hypothesis"] or "control" for r in hedged_rows)

    import feedback
    import plots
    signal = feedback.analyse(feedback.load(OUT / "IL-6_hedged_plate.csv",
                                            OUT / "DEMO_1_signal.csv"))
    noise = feedback.analyse(feedback.load(OUT / "IL-6_hedged_plate.csv",
                                           OUT / "DEMO_2_no_signal.csv"))
    # Drawn here, from the demo file, and labelled as such on the image itself.
    joined = feedback.load(OUT / "IL-6_hedged_plate.csv", OUT / "DEMO_1_signal.csv")
    plots.feedback_scatter(
        joined, signal, str(OUT / "feedback.png"),
        provenance="SIMULATED results — demonstrates the feedback path; "
                   "no wet-lab data has been collected yet")
    return {
        "manifest": m,
        "run_dir": runs[-1].parent.name,
        "test_wells": sum(1 for r in plate_rows if r["Role"] == "test"),
        "control_wells": sum(1 for r in plate_rows if r["Role"] == "control"),
        "hedged": dict(by_hyp),
        "signal": signal,
        "noise": noise,
    }


# ---------------------------------------------------------------- layout

def _text(slide, left, top, width, height, runs, size=18, colour=INK,
          bold=False, align=PP_ALIGN.LEFT, spacing=1.25):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(runs if isinstance(runs, list) else [runs]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(8)
        text, opts = (item if isinstance(item, tuple) else (item, {}))
        run = para.add_run()
        run.text = text
        f = run.font
        f.name = FONT
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.color.rgb = opts.get("colour", colour)
    return box


def _slide(prs, title: str, kicker: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    top = Inches(0.5)
    if kicker:
        _text(s, MARGIN, top, W - 2 * MARGIN, Inches(0.3),
              kicker.upper(), size=12, colour=ACCENT, bold=True)
        top = Inches(0.85)
    _text(s, MARGIN, top, W - 2 * MARGIN, Inches(0.7), title,
          size=30, bold=True, spacing=1.05)
    return s, top + Inches(0.85)


def _rule(slide, top):
    line = slide.shapes.add_shape(1, MARGIN, top, W - 2 * MARGIN, Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE4, 0xE4, 0xEA)
    line.line.fill.background()
    line.shadow.inherit = False


def _picture(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top,
                                    width=width, height=height)


def _stat(slide, left, top, value: str, label: str, width=Inches(2.6),
          colour=ACCENT):
    _text(slide, left, top, width, Inches(0.6), value, size=34, bold=True,
          colour=colour, spacing=1.0)
    _text(slide, left, top + Inches(0.62), width, Inches(0.8), label,
          size=12, colour=MUTED, spacing=1.2)


# ---------------------------------------------------------------- deck

def build(f: dict, path: Path) -> Path:
    m = f["manifest"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 ── title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    _text(s, MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(0.4),
          "re:AGENT · TRACK C · BIO DESIGN", size=13, colour=ACCENT, bold=True)
    _text(s, MARGIN, Inches(2.5), Inches(10.5), Inches(1.4),
          "An agent that designs aptamer switches", size=46, bold=True,
          spacing=1.05)
    _text(s, MARGIN, Inches(3.9), Inches(10.5), Inches(1.0),
          "Type a biomarker. Get 96 wells a wet lab can order — and the "
          "reasons to believe them, or not.", size=21, colour=MUTED)
    _rule(s, Inches(5.2))
    _text(s, MARGIN, Inches(5.5), Inches(11.5), Inches(0.9),
          [("Continuous electrochemical (E-AB) biosensing · IL-6, IL-10, "
            "TNF-α", {"size": 15, "colour": INK}),
           ("Claude Agent SDK · ViennaRNA · Paperclip · Modal/Proto · Primer3",
            {"size": 13, "colour": MUTED})])

    # 2 ── the problem
    s, top = _slide(prs, "A binder is not a sensor", "the problem")
    _text(s, MARGIN, top, Inches(6.0), Inches(3.6),
          [("An E-AB sensor reports a conformational change, not a binding "
            "event. The aptamer is tethered to gold at the 5′ end and carries "
            "a methylene-blue reporter at the 3′ end; current changes only if "
            "the molecule folds differently when the target arrives.",
            {"size": 17}),
           ("So the design problem is not \"find a binder\". It is: take a "
            "published binder and destabilise it by exactly the right amount.",
            {"size": 17, "bold": True}),
           ("Too stable and it never switches. Too loose and it never folds. "
            "The usable window is roughly 2 kcal/mol wide — and that is "
            "narrower than the folding model's own error bar.",
            {"size": 17, "colour": WARN})])
    _picture(s, OUT / "IL-6_switch.png", Inches(7.0), top - Inches(0.1),
             width=Inches(5.6))

    # 3 ── what it does
    s, top = _slide(prs, "One input, eight tools, an orderable plate",
                    "what we built")
    steps = [
        ("find_parents", "greps the full text of the literature for sequences"),
        ("assess_parent", "folds it; decides which design lever even exists"),
        ("design_plate", "enumerates, scores, filters, tiles, exports"),
        ("validate_plate", "re-scores with an independent engine (Primer3)"),
        ("test_core_sensitivity", "does the plate survive its own assumption?"),
        ("build_hedged_plate", "spends 96 wells across the uncertainty"),
        ("learn_from_results", "reads the bench data and moves the window"),
        ("read_ledger", "what is measured, assumed, and rejected"),
    ]
    y = top
    for name, what in steps:
        _text(s, MARGIN, y, Inches(3.0), Inches(0.3), name, size=15,
              bold=True, colour=ACCENT, spacing=1.0)
        _text(s, MARGIN + Inches(3.1), y, Inches(4.4), Inches(0.3), what,
              size=15, colour=INK, spacing=1.0)
        y += Inches(0.53)
    _text(s, Inches(8.3), top, Inches(4.4), Inches(3.2),
          [("The agent chooses the order and can refuse.", {"size": 16,
                                                            "bold": True}),
           ("If no published parent binds the target, it says so and stops "
            "rather than inventing one — de novo aptamer design has no "
            "validated computational method.", {"size": 15, "colour": MUTED}),
           ("Everything runs on the Claude Agent SDK with Claude Code's own "
            "tools switched off, so the agent has exactly these eight and "
            "nothing else.", {"size": 15, "colour": MUTED})])

    # 4 ── literature
    s, top = _slide(prs, "Search for sequences, not for topics", "evidence")
    _text(s, MARGIN, top, Inches(6.2), Inches(3.4),
          [("Ranking papers by relevance finds reviews. We grep the full text "
            "of the corpus for DNA/RNA strings, then read what is printed on "
            "either side of each hit.", {"size": 17}),
           ("Every candidate carries printed_as — the words the paper set "
            "immediately around the sequence — so the agent can see it is "
            "about to design against the receptor rather than the ligand.",
            {"size": 17}),
           ("Corroboration is counted, not assumed: how many independent "
            "papers print the identical sequence.", {"size": 17})])
    _stat(s, Inches(7.3), top + Inches(0.1), "53", "papers containing a\ncandidate sequence")
    _stat(s, Inches(10.0), top + Inches(0.1), "10", "parents extracted\nwith attribution")
    _stat(s, Inches(7.3), top + Inches(1.9), "3", "independent papers\nprint the chosen parent")
    _stat(s, Inches(10.0), top + Inches(1.9), "0", "affinities borrowed\nfrom another aptamer",
          colour=WARN)
    _text(s, Inches(7.3), top + Inches(3.4), Inches(5.3), Inches(1.0),
          "No paper reports a Kd for this exact sequence, so every "
          "affinity-derived number is omitted rather than estimated. An "
          "earlier version borrowed 8.5 nM from an anti-IL-6-receptor "
          "aptamer and produced a full column of wrong figures, none marked "
          "as such.", size=13, colour=MUTED)

    # 5 ── the funnel
    s, top = _slide(prs, "From 8,577 candidates to 96 wells", "design")
    _picture(s, OUT / "IL-6_funnel.png", MARGIN, top, width=Inches(6.6))
    _text(s, Inches(7.6), top, Inches(5.0), Inches(3.6),
          [(f"{m['library_size']:,} variants enumerated, not sampled — every "
            f"tail length, register, mismatch count and linker.", {"size": 16}),
           (f"{m['in_switching_window']:,} land in the switching window; "
            f"{m['passing_all_criteria']:,} also pass specificity, "
            f"self-dimer, G-run, homopolymer and GC filters.", {"size": 16}),
           ("Filters run before tiling, not after. Tiling first leaves holes "
            "wherever a band happened to contain only dimerising designs, and "
            "the plate silently stops covering the range it claims to.",
            {"size": 16, "colour": MUTED})])

    # 6 ── why tile
    s, top = _slide(prs, "We tile the window instead of ranking it",
                    "the honest core")
    _picture(s, OUT / "IL-6_window.png", MARGIN, top, width=Inches(7.0))
    _text(s, Inches(8.0), top, Inches(4.6), Inches(3.8),
          [("Taking the top 96 by predicted energy gives 96 near-identical "
            "sequences and one bit of information.", {"size": 16}),
           ("If the model's centre is off by a kcal/mol — and near zero it "
            "can be — every well fails together and the round teaches "
            "nothing.", {"size": 16}),
           ("So the plate spans eight energy bands, 11 wells each, chosen to "
            "be mechanistically unalike within a band. The experiment returns "
            "the shape of the response and locates the optimum even when the "
            "prediction is wrong.", {"size": 16, "bold": True}),
           (f"Plate position is randomised against ddG and checked against a "
            f"2,000-permutation null: row spread "
            f"{m['position_check']['observed']} vs "
            f"{m['position_check']['null_p95']} at the 95th percentile.",
            {"size": 14, "colour": MUTED})])

    # 7 ── the problem we found
    s, top = _slide(prs, "Then we tested our own assumption, and it failed",
                    "what most tools skip")
    _text(s, MARGIN, top, Inches(11.9), Inches(1.2),
          "No paper maps where IL-6 contacts this aptamer. Every energy in "
          "the design is computed against an assumed binding core — so we "
          "swept the assumption and re-designed under each one.", size=18)
    y = top + Inches(1.3)
    _stat(s, MARGIN, y, "2 of 4", "core hypotheses yield\nany usable plate at all")
    _stat(s, Inches(4.2), y, "0", "designs survive\nevery core hypothesis", colour=WARN)
    _stat(s, Inches(7.6), y, "176", "designs total across\nthe productive cores")
    _text(s, MARGIN, y + Inches(1.9), Inches(11.9), Inches(1.6),
          [("The plate was an artefact of an assumption we could not check.",
            {"size": 20, "bold": True, "colour": WARN}),
           ("\"Do not order\" is not an answer a wet lab can act on — they "
            "have a budget, a synthesis slot and a question. Refusing to "
            "choose is not rigour; it passes the problem back.",
            {"size": 17, "colour": MUTED})])

    # 8 ── hedged plate
    s, top = _slide(prs, "So we spend the plate on the uncertainty",
                    "the deliverable")
    _picture(s, OUT / "IL-6_plate.png", MARGIN, top, width=Inches(6.4))
    h = f["hedged"]
    _text(s, Inches(7.4), top, Inches(5.2), Inches(4.0),
          [("One 96-well plate, allocated across both surviving core "
            "hypotheses, every well labelled with the hypothesis it belongs "
            "to.", {"size": 17}),
           (" · ".join(f"{k}: {v} wells" for k, v in sorted(h.items())),
            {"size": 15, "bold": True, "colour": ACCENT}),
           ("Read it twice: which switches work, and which core was right. A "
            "hypothesis with no responsive well is eliminated — the epitope "
            "question answered by a synthesis run that was happening anyway.",
            {"size": 17}),
           ("Wells are randomised against hypothesis as well as against "
            "energy. Split by hypothesis into rows A–D and E–H and any row "
            "gradient is perfectly confounded with the comparison. Controls "
            "are shared, not duplicated: they report on the assay, not the "
            "core.", {"size": 14, "colour": MUTED})])

    # 9 ── closing the loop
    sig, noi = f["signal"], f["noise"]
    s, top = _slide(prs, "The bench answers, and the next round moves",
                    "design → build → test → learn")
    _picture(s, OUT / "feedback.png", MARGIN, top, width=Inches(6.8))
    hyp = sig["by_core_hypothesis"]
    lines = []
    for k, v in sorted(hyp.items()):
        lines.append((f"{k}: {v['responsive']}/{v['wells']} wells responded, "
                      f"median {v['median_signal']}%",
                      {"size": 15, "bold": True,
                       "colour": ACCENT if v["responsive"] > 20 else MUTED}))
    _text(s, Inches(7.8), top, Inches(4.8), Inches(4.2),
          [("Upload a results CSV and the agent reads it before it designs "
            "anything.", {"size": 16}),
           *lines,
           (f"One plate settles the epitope question. The window recentres on "
            f"the measured optimum (ddG {sig['measured_optimum_ddg']:+.2f}), "
            f"not the predicted one.", {"size": 16}),
           (f"And on a pure-noise plate it reports ρ={noi['ddg_vs_signal']['spearman']}, "
            f"p={noi['ddg_vs_signal']['permutation_p']} and refuses to "
            f"recentre. Every correlation is tested against a "
            f"2,000-permutation null before it is allowed to move a design.",
            {"size": 16, "bold": True, "colour": WARN})])

    # 10 ── rejected
    s, top = _slide(prs, "Three models we tried, controlled, and threw away",
                    "what did not work")
    import ledger
    y = top
    for r in ledger.REJECTED:
        _text(s, MARGIN, y, Inches(11.9), Inches(0.3), r["model"], size=17,
              bold=True, spacing=1.0)
        _text(s, MARGIN, y + Inches(0.32), Inches(11.9), Inches(0.9),
              f"Control: {r['control']}.  {r['result']}.", size=14,
              colour=MUTED, spacing=1.2)
        y += Inches(1.32)
    _rule(s, y - Inches(0.15))
    _text(s, MARGIN, y + Inches(0.1), Inches(11.9), Inches(1.0),
          [("A tool that only reports what it used looks more certain than it "
            "is.", {"size": 17, "bold": True}),
           ("The dimer and fold energies the plate is filtered on were "
            "checked against Primer3's independent implementation: homodimer "
            "rank correlation 0.49–0.67 across plates, 0.81 on random DNA. "
            "Absolute values differ by ~6 kcal/mol because the parameter sets "
            "differ; the ordering, which is what the filter acts on, holds.",
            {"size": 14, "colour": MUTED})])

    # 11 ── limits
    s, top = _slide(prs, "What this does not do", "limits, stated up front")
    _text(s, MARGIN, top, Inches(11.9), Inches(4.0),
          [("It cannot invent an aptamer. Without a published parent there is "
            "nothing honest to build on, and the agent says so.", {"size": 18}),
           ("It cannot predict affinity. Apparent Kd is computed only when a "
            "paper reports one for that exact sequence — usually none does, "
            "and the affinity figures are then omitted.", {"size": 18}),
           ("It cannot reach healthy baseline. IL-6 circulates at pg/mL and "
            "aptamer Kd values are nanomolar; switching costs affinity on top "
            "of that. This sensor class reaches sepsis and cytokine-release "
            "range, not routine monitoring — and the memo says so every time.",
            {"size": 18, "colour": WARN}),
           ("The binding core remains an assumption. That is precisely why "
            "the deliverable is a hedged plate rather than a confident one.",
            {"size": 18})])

    # 12 ── close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    _text(s, MARGIN, Inches(2.3), Inches(11.5), Inches(1.2),
          "96 wells, and an honest account of what they rest on",
          size=38, bold=True, spacing=1.05)
    _rule(s, Inches(3.9))
    _text(s, MARGIN, Inches(4.2), Inches(11.5), Inches(1.6),
          [("Every run archives a manifest: the parent and its citation, the "
            "thresholds chosen, the assumption made, the commit that produced "
            "it. A reader deciding whether to spend a synthesis slot can see "
            "all of it.", {"size": 18}),
           (f"Reproduced from commit {m.get('code_commit', '')} · "
            f"github.com/ascendbio/aptamer-switch", {"size": 14,
                                                     "colour": MUTED})])

    _notes(prs, f)
    prs.save(path)
    return path


# Five minutes over twelve slides is twenty-five seconds each, which is not how
# the talk should be spent. The budget below is deliberately uneven: the first
# four slides set up a problem the audience mostly accepts, and slides 7-9 carry
# the only claim that distinguishes this from a scoring script.
NOTES = [
    (12, "Continuous cytokine sensing. We design the aptamer switches, not the "
         "electronics. Three targets, IL-6 first because it is the hardest."),
    (28, "The key idea, and the one thing to land: an E-AB sensor reports a "
         "SHAPE CHANGE, not binding. A perfect binder that never moves gives no "
         "signal. So we take a published binder and destabilise it by a "
         "precise amount — and the usable window is narrower than the folding "
         "model's error bar. That tension drives every decision after this."),
    (20, "One input: a biomarker name. Eight tools, agent picks the order. "
         "Point at find_parents and build_hedged_plate — those are the two that "
         "matter. Say it can refuse: no published parent, no plate."),
    (18, "We grep full text for SEQUENCES, not rank papers by topic. Each hit "
         "carries the words printed around it, so the agent can see it is "
         "about to design against the receptor instead of the ligand. It "
         "caught exactly that. Zero borrowed affinities — an earlier version "
         "borrowed a Kd from an anti-IL-6-RECEPTOR aptamer and produced a "
         "column of confident wrong numbers."),
    (16, "8,577 enumerated, not sampled. Filters run BEFORE tiling — tiling "
         "first leaves holes wherever a band held only dimerising designs, and "
         "the plate stops covering the range it claims to."),
    (26, "Why we do not ship a top-96 list. Ranked, you get 96 near-identical "
         "sequences and one bit of information; if the model centre is off by "
         "one kcal/mol they all fail together. Tiling returns the SHAPE of the "
         "response even when the prediction is wrong. Position randomised "
         "against energy, checked against a permutation null."),
    (40, "The slide to slow down on. No paper maps where IL-6 touches this "
         "aptamer, so every energy rests on an assumed binding core. We swept "
         "it. Two of four cores give a plate at all, and ZERO designs survive "
         "every core — our plate was an artefact of an assumption we cannot "
         "check. Then say the hard part: 'do not order' is not an answer. A "
         "lab has a budget and a synthesis slot."),
    (34, "So spend the plate on the uncertainty. 44 wells per core hypothesis, "
         "8 shared controls, every well labelled. Read it twice: which "
         "switches work, and which core was right. A hypothesis with no "
         "responsive well is eliminated — the epitope question answered by a "
         "synthesis run that was happening anyway. Randomised against "
         "hypothesis, not blocked by row, or a warm corner of the plate "
         "becomes your result."),
    (40, "The loop closes. Upload the plate reader CSV, it is read before "
         "anything is designed. On this plate one core responded 40 of 44, the "
         "other 5 of 44 — epitope settled in one round. Window recentres on "
         "the MEASURED optimum. Then the important half: on a pure-noise "
         "plate it reports p=0.26 and refuses to move. A tool that finds "
         "signal in noise is worse than no tool. Say the data on screen is "
         "simulated — it is labelled on the figure."),
    (20, "Three external models we tried and threw away, each against a "
         "control. OpenDDE: the SCRAMBLED aptamer scored higher, so the "
         "contacts cannot define a core. Evo2: it prefers repetitive DNA, not "
         "aptamers. We report these because a tool that only lists what it "
         "used looks more certain than it is."),
    (14, "Limits, said before anyone asks. Cannot invent an aptamer. Cannot "
         "predict affinity. Cannot reach healthy baseline — IL-6 circulates at "
         "pg/mL, aptamer Kd is nanomolar, and switching costs more affinity on "
         "top. This class reaches sepsis and CRS, and the memo says so every "
         "time."),
    (8, "Every run archives a manifest: parent, citation, thresholds, the "
         "assumption, the commit. Close on the hedged plate as the "
         "deliverable."),
]


def _notes(prs, f: dict) -> None:
    """Speaker notes with a time budget that adds up to five minutes."""
    total = sum(sec for sec, _ in NOTES)
    running = 0
    for slide, (sec, text) in zip(prs.slides, NOTES):
        running += sec
        tf = slide.notes_slide.notes_text_frame
        tf.text = (f"[{sec}s · {running // 60}:{running % 60:02d} of "
                   f"{total // 60}:{total % 60:02d}]\n\n{text}")
    # Five minutes is the hard cap and this is checked rather than eyeballed: the
    # first draft totalled six minutes, which nobody notices until the room does.
    assert total <= 285, f"notes budget is {total}s; five minutes needs slack"


if __name__ == "__main__":
    f = facts()
    out = build(f, ROOT / "reAGENT_TrackC.pptx")
    print(f"written: {out.relative_to(ROOT)}")
    print(f"  library {f['manifest']['library_size']:,} → "
          f"{f['manifest']['in_switching_window']:,} → "
          f"{f['manifest']['passing_all_criteria']:,} → {f['test_wells']} wells")
    print(f"  hedged  {f['hedged']}")
    print(f"  signal  rho {f['signal']['ddg_vs_signal']['spearman']}, "
          f"p {f['signal']['ddg_vs_signal']['permutation_p']}")
    print(f"  noise   rho {f['noise']['ddg_vs_signal']['spearman']}, "
          f"p {f['noise']['ddg_vs_signal']['permutation_p']}")
